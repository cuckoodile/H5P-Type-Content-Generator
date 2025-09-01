from pptx import Presentation
import os

def ppt_reader(file_path):
    if not os.path.exists(file_path):
        print(f"[ERROR] File not found: {file_path}")
        return ""

    presentation = Presentation(file_path)
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    output_lines = []

    for i, slide in enumerate(presentation.slides):
        output_lines.append(f"\n--- Slide {i + 1} ---")
        for shape in slide.shapes:
            numbered_index = 1  # Reset numbering per shape

            # Handle text frames
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs: # type: ignore
                    try:
                        text = paragraph.text.strip()
                        if not text:
                            continue

                        bullet_type = None
                        p_xml = paragraph._element

                        if p_xml.find(".//a:buChar", namespaces=nsmap) is not None:
                            bullet_type = "Bullet"
                        elif p_xml.find(".//a:buAutoNum", namespaces=nsmap) is not None:
                            bullet_type = "Numbered"

                        if bullet_type == "Bullet":
                            output_lines.append(f"* {text}")
                        elif bullet_type == "Numbered":
                            output_lines.append(f"{numbered_index}. {text}")
                            numbered_index += 1
                        else:
                            output_lines.append(text)

                    except Exception as e:
                        output_lines.append(f"[WARN] Failed to parse paragraph: {e}")

            # Handle tables
            elif hasattr(shape, "has_table") and shape.has_table:
                try:
                    output_lines.append("[Table]")
                    for row in shape.table.rows: # type: ignore
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            output_lines.append(" | ".join(row_text))
                except Exception as e:
                    output_lines.append(f"[WARN] Failed to parse table: {e}")

    return "\n".join(output_lines)

if __name__ == "__main__":
    pptx_file = "../references/Managing-Work-Goal-Development-Ch-1.pptx"
    extracted_text = ppt_reader(pptx_file)
    print(extracted_text)
